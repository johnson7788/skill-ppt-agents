"""企业 Agent 自定义工具。

工具：
- todo:           任务规划/进度跟踪，状态存于 tool_context.state（单次会话内有效）
- terminal:       在服务器上执行 shell 命令（高权限，请在受信部署边界内使用）
- vision_analyze: 图片内容分析 / OCR（单独的视觉模型）
- clarify:        人在回路澄清提问（LongRunningFunctionTool，需用户回答后续跑）

注册方式见 app/agent.py：普通函数用 FunctionTool 包装，clarify 直接用
已包装好的 clarify_tool 加入 tools=[...]。
"""

from __future__ import annotations

import base64
import os
import pathlib
import shutil
import subprocess

from google.adk.tools import LongRunningFunctionTool, ToolContext

# 上传目录（与 agent.py 一致；此处本地定义以避免循环导入）
UPLOADS_DIR = pathlib.Path(__file__).parent.parent / "uploads"

# ---------------------------------------------------------------------------
# todo — 任务规划
# ---------------------------------------------------------------------------
_TODO_KEY = "_todo_list"
_VALID_STATUS = {"pending", "in_progress", "completed", "cancelled"}
_MAX_ITEMS = 256
_MAX_CONTENT = 4000


def todo(tool_context: ToolContext, description: str = "", todos: list[dict] | None = None) -> dict:
    """管理本次任务的待办清单，用于拆解复杂任务并跟踪进度。

    用法：
    - description: 操作目的
    - 传入 todos 写入/覆盖整张清单；不传则只读取当前清单。
    - 每个待办项为 {"id": str, "content": str, "status": "..."}，
      status 取值：pending / in_progress / completed / cancelled。
    - 列表顺序即优先级。复杂任务（>=3 步）开始时先写清单，
      每完成一步立即把对应项更新为 completed，并把下一项设为 in_progress。

    返回当前完整清单 {"todos": [...]}。
    """
    items = list(tool_context.state.get(_TODO_KEY, []))
    if todos is not None:
        cleaned: list[dict] = []
        for t in todos[:_MAX_ITEMS]:
            if not isinstance(t, dict):
                continue
            status = str(t.get("status", "pending")).strip().lower()
            cleaned.append({
                "id": str(t.get("id", "")).strip() or str(len(cleaned) + 1),
                "content": str(t.get("content", "")).strip()[:_MAX_CONTENT],
                "status": status if status in _VALID_STATUS else "pending",
            })
        items = cleaned
        tool_context.state[_TODO_KEY] = items
    return {"todos": items}


# ---------------------------------------------------------------------------
# terminal — shell 命令执行
# ---------------------------------------------------------------------------
_TERMINAL_MAX_STDOUT = 8000
_TERMINAL_MAX_STDERR = 2000


async def terminal(description: str = "", command: str = "", tool_context: ToolContext | None = None, timeout: int = 60) -> dict:
    """执行 shell 命令并返回输出。

    适用于查看文件、运行构建/检索脚本、查询系统信息等。
    命令在当前租户的隔离沙箱中执行（多租户互不干扰）；未启用沙箱时在服务器本地执行。

    参数：
    - description: 操作目的
    - command: 要执行的 shell 命令。
    - timeout: 超时秒数（默认 60）。

    返回 {"stdout", "stderr", "returncode"}；超时返回 {"error"}。
    """
    if not command or not command.strip():
        return {"error": "command 不能为空"}

    from app import sandbox as sbx

    if sbx.enabled():
        key = str(tool_context.state.get("_sbkey") or tool_context.user_id or "default")
        try:
            return await sbx.run_async(key, command, timeout)
        except Exception as e:  # noqa: BLE001 — 工具边界，统一回报给模型
            return {"error": f"沙箱执行失败: {e}"}

    try:
        r = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
        )
        return {
            "stdout": r.stdout[-_TERMINAL_MAX_STDOUT:],
            "stderr": r.stderr[-_TERMINAL_MAX_STDERR:],
            "returncode": r.returncode,
        }
    except subprocess.TimeoutExpired:
        return {"error": f"命令超时（>{timeout}s）"}
    except Exception as e:  # noqa: BLE001 — 工具边界，统一回报给模型
        return {"error": str(e)}


# ---------------------------------------------------------------------------
# vision_analyze — 图片内容分析 / OCR
# ---------------------------------------------------------------------------
# DeepSeek 无视觉能力，这里单独用一个视觉模型（默认 qwen-vl-max，OpenAI 兼容端点）。
_VISION_MODEL = os.environ.get("VISION_MODEL", "openai/qwen-vl-max")
_VISION_KEY = os.environ.get("VISION_API_KEY")
_VISION_BASE = os.environ.get(
    "VISION_API_BASE", "https://dashscope.aliyuncs.com/compatible-mode/v1"
)
_VISION_MAX_BYTES = 10 * 1024 * 1024  # 单图 10MB 上限
_IMAGE_MIME = {
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".webp": "image/webp", ".gif": "image/gif", ".bmp": "image/bmp",
}


def vision_analyze(description: str = "", image: str = "", prompt: str = "描述这张图片的内容，并识别其中的文字") -> dict:
    """分析图片内容或对图片做 OCR 文字识别。

    适用于用户上传了图片（如截图、图表、扫描件、含文字的图片）需要理解或提取文字的场景。
    image 可以是 http(s) 图片 URL，或用户已上传的图片文件名（位于 uploads 目录）。
    prompt 描述你想从图片中获取什么信息（如"提取图中表格数据"、"这是什么模型结构图"）。

    参数：
    - description: 操作目的
    - image: 图片 URL 或已上传的文件名。
    - prompt: 分析提示词（可选）。

    返回 {"analysis": str}，失败返回 {"error": str}。
    """
    if not _VISION_KEY:
        return {"error": "未配置 VISION_API_KEY，无法使用图片分析。请在 .env 中设置视觉模型。"}

    try:
        if image.startswith(("http://", "https://")):
            import httpx
            resp = httpx.get(image, timeout=30, follow_redirects=True)
            resp.raise_for_status()
            data = resp.content
            mime = resp.headers.get("content-type", "image/jpeg").split(";")[0]
            if not mime.startswith("image/"):
                mime = "image/jpeg"
        else:
            p = next(UPLOADS_DIR.rglob(image), None)
            if p is None or not p.is_file():
                return {"error": f"未找到图片文件: {image}（请确认已上传）"}
            data = p.read_bytes()
            mime = _IMAGE_MIME.get(p.suffix.lower(), "image/jpeg")

        if len(data) > _VISION_MAX_BYTES:
            return {"error": f"图片过大（{len(data) // 1024 // 1024}MB > 10MB），请压缩后重试。"}

        b64 = base64.b64encode(data).decode()
        import litellm
        completion = litellm.completion(
            model=_VISION_MODEL,
            api_key=_VISION_KEY,
            api_base=_VISION_BASE,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                ],
            }],
        )
        return {"analysis": completion.choices[0].message.content}
    except Exception as e:  # noqa: BLE001 — 工具边界，统一回报给模型
        return {"error": str(e)}


# ---------------------------------------------------------------------------
# clarify — 人在回路澄清提问
# ---------------------------------------------------------------------------
def clarify(question: str, choices: list[str] | None = None) -> dict:
    """当用户需求不明确、存在多种可能理解，或缺少关键信息无法继续时，向用户提问澄清。

    用法：
    - question: 你想向用户确认的问题，应具体、简洁。
    - choices: 可选项列表（最多 4 项）。提供选项能让用户更快回答；
      若是开放式问题可不传。

    这是一个长时运行工具：调用后会暂停并等待用户回答，用户答复后你将
    收到答案并据此继续。**遇到任何不确定的信息都必须使用此工具向用户确认，禁止自行猜测用户意图。**
    """
    return {
        "status": "pending",
        "question": question,
        "choices": list(choices or [])[:4],
    }


# 包装为长时运行工具：调用后挂起，等待用户通过 /chat/answer 回灌答复。
clarify_tool = LongRunningFunctionTool(func=clarify)


# ---------------------------------------------------------------------------
# upload_to_sandbox — 上传文件到沙箱
# ---------------------------------------------------------------------------
def upload_to_sandbox(
    sandbox_path: str,
    content: str,
    tool_context: ToolContext,
    mode: int = 644,
) -> dict:
    """上传文件到当前租户的沙箱内。

    适用于需要在沙箱中预置脚本、配置文件或数据文件的场景。
    文件写入沙箱后，可通过 terminal 工具在沙箱中访问。

    参数：
    - sandbox_path: 沙箱内的目标路径（如 /scripts/my_tool.py）。
    - content: 文件内容（纯文本）。
    - mode: 文件权限（默认 644）。

    返回 {"success": true} 或 {"error": str}。
    """
    from app import sandbox as sbx

    if not sbx.enabled():
        return {"error": "沙箱未启用，无法上传文件"}

    key = str(tool_context.state.get("_sbkey") or tool_context.user_id or "default")
    try:
        sbx.write_file_sync(key, sandbox_path, content, mode)
        return {"success": True, "path": sandbox_path}
    except Exception as e:
        return {"error": f"上传文件到沙箱失败: {e}"}


# ---------------------------------------------------------------------------
# ensure_sandbox_skills — 同步指定 skill 脚本到沙箱
# ---------------------------------------------------------------------------
def ensure_sandbox_skills(skill_name: str, tool_context: ToolContext) -> dict:
    """确保当前租户的沙箱内已有指定 skill 的脚本文件。

    首次使用某个 skill 前调用一次即可。会比对沙箱现有文件，仅同步缺失的文件。
    同步后可通过 terminal 在沙箱中直接运行该 skill 的脚本。

    参数：
    - skill_name: 技能目录名（如 "bingsearch"、"arxiv-paper-search"）。

    返回 {"success": true} 或 {"error": str}。
    """
    from app import sandbox as sbx

    if not sbx.enabled():
        return {"error": "沙箱未启用，无法同步 skill"}

    if not skill_name or not skill_name.strip():
        return {"error": "skill_name 不能为空"}

    key = str(tool_context.state.get("_sbkey") or tool_context.user_id or "default")
    try:
        sbx.ensure_skills_sync(key, skill_name.strip())
        return {"success": True, "message": f"skill '{skill_name}' 已同步到沙箱"}
    except FileNotFoundError:
        return {"error": f"本地不存在 skill '{skill_name}'"}
    except Exception as e:
        return {"error": f"同步 skill 到沙箱失败: {e}"}


# ---------------------------------------------------------------------------
# sync_upload_to_sandbox — 将用户已上传的文件同步到沙箱
# ---------------------------------------------------------------------------
def sync_upload_to_sandbox(
    filename: str,
    tool_context: ToolContext,
    sandbox_path: str | None = None,
) -> dict:
    """将用户已上传的文件从服务器同步到当前租户的沙箱内。

    适用于需要将用户上传的数据文件传入沙箱进行处理的场景。
    调用前应先用 list_uploads 查看已上传文件列表。

    参数：
    - filename: 已上传的文件名（如 "data.csv"、"report.pdf"）。
    - sandbox_path: 沙箱内的目标路径（默认 /uploads/<filename>）。
      注意：沙箱内路径必须以 / 开头，如 "/workspace/data.csv"。

    返回 {"success": true, "sandbox_path": "..."} 或 {"error": str}。
    """
    from app import sandbox as sbx

    if not sbx.enabled():
        return {"error": "沙箱未启用，无法同步文件"}

    user_id = str(tool_context.state.get("_sbkey") or tool_context.user_id or "default")
    user_dir = UPLOADS_DIR / user_id

    if not user_dir.is_dir():
        return {"error": f"用户 {user_id} 的上传目录不存在"}

    safe_name = pathlib.Path(filename).name
    src = user_dir / safe_name

    if not src.is_file():
        return {"error": f"文件 '{filename}' 不存在，请先上传"}

    try:
        resolved = src.resolve()
        user_dir_resolved = user_dir.resolve()
        if not str(resolved).startswith(str(user_dir_resolved)):
            return {"error": "权限拒绝：禁止访问其他用户的上传目录"}
    except (OSError, ValueError):
        return {"error": "路径解析失败"}

    target = sandbox_path or f"/uploads/{safe_name}"
    try:
        data = src.read_bytes()
        sbx.write_file_sync(user_id, target, data, mode=644)
        return {"success": True, "sandbox_path": target, "size": len(data)}
    except Exception as e:
        return {"error": f"同步文件到沙箱失败: {e}"}


# ---------------------------------------------------------------------------
# generate_ppt — 图片型 PPT 生成（qwen-image 逐页出图 + python-pptx 组装）
# ---------------------------------------------------------------------------
# 每页是一整张 16:9 生成图，最后组装成 .pptx。图片后端用阿里 DashScope
# qwen-image（异步任务）。生成+组装是纯确定性活，放在后端本地跑（后端已带
# requests + python-pptx，且能直接写 uploads/ 供下载），不进沙箱——沙箱镜像
# 没有 python-pptx，也没有把产物拉回宿主机的通道。
_QWEN_IMAGE_API = "https://dashscope.aliyuncs.com/api/v1/services/aigc/text2image/image-synthesis"
_QWEN_IMAGE_TASK = "https://dashscope.aliyuncs.com/api/v1/tasks/"
_PPT_MAX_SLIDES = 20
# qwen-image 无负向提示词：写"不要浏览器"反而会把浏览器画出来，所以只写正向约束。
_PPT_QUALITY_RULES = (
    "16:9 full-slide PowerPoint image. No watermark, no color codes or hex text, "
    "no page numbers, no scaffolding labels. No browser chrome, no application "
    "toolbar or menu bar, no window frame, no address bar. Render every Chinese character crisply "
    "and correctly in a clean Chinese sans-serif; keep any latin terms spelled exactly."
)
# 预设风格 → 英文视觉描述（注入每页出图提示词）。风格取自 codex-ppt references。
_STYLE_PRESETS: dict[str, str] = {
    "科研答辩风": "formal Chinese academic research-defense style, clean white background, deep academic blue structure, research-blue accents, pale-blue fills, one formal-red emphasis, precise alignment",
    "麦肯锡风格": "McKinsey consulting style, clean white background, navy and steel-blue palette, minimal, strong grid, thin horizontal dividers, one accent color, data-driven",
    "清爽专业风": "clean professional style, white background, soft blue and light-gray palette, generous whitespace, rounded cards, modern sans-serif",
    "数据仪表盘风": "data dashboard style, deep navy background, cyan and teal accents, KPI cards, charts and gridlines, modern",
    "党政红风格": "Chinese government report style, warm red and gold palette, cream background, solemn formal banner headers",
    "教学课件风": "classroom courseware style, bright friendly blue and orange palette, clear headings, simple diagrams, approachable",
    "温暖手工风": "warm handcraft style, cream paper background, warm earth tones, soft rounded shapes, cozy",
    "手绘白板风": "hand-drawn whiteboard style, white background, black marker outlines, sketchy hand-drawn diagrams, a few accent colors",
    "手绘技术解释风": "hand-drawn technical explainer style, off-white background, ink line-art diagrams, muted accent colors, annotated",
    "电子墨水杂志风": "e-ink magazine style, warm paper background, monochrome black-and-cream with one accent, editorial serif headings, calm",
    "创意杂志风": "creative magazine editorial style, bold color blocks, large expressive typography, asymmetric layout, vivid",
    "复古扁平插画风": "retro flat illustration style, muted vintage palette, simple flat shapes, textured background, mid-century",
}
_DEFAULT_STYLE = "科研答辩风"


def _qwen_gen_image(prompt: str, out_path: str, key: str, size: str = "1664*928") -> None:
    """调 qwen-image 异步任务出一张图并下载到 out_path。失败抛异常。"""
    import json
    import time
    import urllib.request

    def _req(url, body=None, extra=None):
        data = json.dumps(body).encode() if body is not None else None
        r = urllib.request.Request(url, data=data, method="POST" if body else "GET")
        r.add_header("Authorization", f"Bearer {key}")
        if body is not None:
            r.add_header("Content-Type", "application/json")
        for k, v in (extra or {}).items():
            r.add_header(k, v)
        with urllib.request.urlopen(r, timeout=120) as resp:
            return json.loads(resp.read())

    sub = _req(
        _QWEN_IMAGE_API,
        {"model": "qwen-image", "input": {"prompt": prompt},
         "parameters": {"size": size, "n": 1, "prompt_extend": False, "watermark": False}},
        {"X-DashScope-Async": "enable"},
    )
    tid = sub["output"]["task_id"]
    for _ in range(60):
        time.sleep(5)
        res = _req(_QWEN_IMAGE_TASK + tid)
        st = res["output"]["task_status"]
        if st == "SUCCEEDED":
            url = res["output"]["results"][0]["url"]
            urllib.request.urlretrieve(url, out_path)
            return
        if st in ("FAILED", "UNKNOWN"):
            raise RuntimeError(f"qwen-image 任务 {st}: {res.get('output')}")
    raise TimeoutError("qwen-image 任务超时")


def _assemble_pptx(image_paths: list[str], notes: list[str], out_path: str) -> None:
    """把每张图铺满一页 16:9 幻灯片，写入演讲备注，保存 .pptx。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]
    for i, img in enumerate(image_paths):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
        note = notes[i] if i < len(notes) else ""
        if note:
            slide.notes_slide.notes_text_frame.text = note
    prs.save(out_path)


def _run_ppt_pipeline(user_id: str, title: str, style: str, slides: list[dict]) -> dict:
    """阻塞式：逐页出图 + 组装。放在线程里跑，避免阻塞事件循环。"""
    key = os.environ.get("DASHSCOPE_API_KEY")
    if not key:
        return {"error": "未配置 DASHSCOPE_API_KEY，无法生成 PPT。请在 .env 中设置。"}
    style_desc = _STYLE_PRESETS.get(style, _STYLE_PRESETS[_DEFAULT_STYLE])
    safe_title = "".join(c for c in title if c not in '\\/:*?"<>|').strip() or "deck"
    user_dir = UPLOADS_DIR / user_id
    work_dir = user_dir / "_ppt" / safe_title
    work_dir.mkdir(parents=True, exist_ok=True)

    image_paths: list[str] = []
    notes: list[str] = []
    for idx, s in enumerate(slides, 1):
        body = (s.get("prompt") or "").strip()
        if not body:
            return {"error": f"第 {idx} 页缺少 prompt"}
        full = f"{_PPT_QUALITY_RULES} Visual style: {style_desc}. {body}"
        out = work_dir / f"slide_{idx:02d}.png"
        try:
            _qwen_gen_image(full, str(out), key)
        except Exception as e:  # noqa: BLE001 — 出图边界，回报给模型
            return {"error": f"第 {idx} 页出图失败: {e}"}
        image_paths.append(str(out))
        notes.append((s.get("notes") or "").strip())

    pptx_name = f"{safe_title}.pptx"
    pptx_path = user_dir / pptx_name
    try:
        _assemble_pptx(image_paths, notes, str(pptx_path))
    except Exception as e:  # noqa: BLE001
        return {"error": f"组装 PPT 失败: {e}"}

    from urllib.parse import quote
    return {
        "success": True,
        "slides": len(image_paths),
        "file": pptx_name,
        "download_url": f"/download?user_id={quote(user_id)}&file={quote(pptx_name)}",
    }


async def generate_ppt(
    title: str,
    slides: list[dict],
    tool_context: ToolContext,
    style: str = _DEFAULT_STYLE,
    description: str = "",
) -> dict:
    """生成图片型 PPT：每页一整张 16:9 生成图，组装为 .pptx 供下载。

    你（模型）负责创意部分——先规划提纲，再为每一页写好出图提示词；本工具负责
    机械部分——逐页调用图片模型出图并组装。调用前请把每一页的画面/版式想清楚。

    参数：
    - title: 演示文稿标题（也作为文件名）。
    - slides: 每页一个 dict：{"prompt": 该页出图提示词, "notes": 可选演讲备注}。
      prompt 用英文骨架描述版式、用中文写要显示的标题/正文文字，文字务必精简、
      每个字都要正确；不要写十六进制色号、不要写"Card 1/2"之类脚手架词。
    - style: 预设风格名，取值之一：科研答辩风、麦肯锡风格、清爽专业风、数据仪表盘风、
      党政红风格、教学课件风、温暖手工风、手绘白板风、手绘技术解释风、
      电子墨水杂志风、创意杂志风、复古扁平插画风。默认科研答辩风。
    - description: 操作目的（展示用）。

    返回 {"success", "slides", "file", "download_url"}；失败返回 {"error"}。
    生成后请把 download_url 以 markdown 链接形式给用户，让其点击下载。
    出图较慢（每页约需 30-60 秒），请一次把整套页面传入。
    """
    import asyncio

    if not isinstance(slides, list) or not slides:
        return {"error": "slides 不能为空"}
    if len(slides) > _PPT_MAX_SLIDES:
        return {"error": f"页数过多（{len(slides)} > {_PPT_MAX_SLIDES}），请精简"}
    user_id = str(
        (tool_context.state.get("_sbkey") if tool_context else None)
        or (tool_context.user_id if tool_context else None)
        or "default_user"
    )
    return await asyncio.to_thread(_run_ppt_pipeline, user_id, title, style, slides)
